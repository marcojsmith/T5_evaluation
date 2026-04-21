# -*- coding: utf-8 -*-
from pptx import Presentation
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

prs = Presentation('presentation.pptx')

results = []

for i, slide in enumerate(prs.slides):
    slide_text = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            slide_text.append(shape.text.strip())
    results.append((i+1, '\n'.join(slide_text)))

for slide_num, content in results:
    print(f"SLIDE {slide_num}")
    print(content)
    print("---END---")