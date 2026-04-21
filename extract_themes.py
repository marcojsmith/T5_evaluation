# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import sys
import io
import json

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

prs = Presentation('presentation.pptx')

for i, slide in enumerate(prs.slides):
    slide_num = i + 1
    backgrounds = []
    for shape in slide.shapes:
        if shape.shape_type == 1:  # MSO_SHAPE_TYPE.RECTANGLE
            if hasattr(shape, 'fill'):
                fill = shape.fill
                if fill.type == 1:  # solid fill
                    bg = fill.fore_color
                    if hasattr(bg, 'rgb'):
                        backgrounds.append(bg.rgb)
    
    # Check slide background
    try:
        bg = slide.background
        if bg.fill.type == 1:
            f = bg.fill.fore_color
            if hasattr(f, 'rgb'):
                backgrounds.append(f.rgb)
    except:
        pass
    
    # Determine theme based on background
    theme = "unknown"
    for bg in backgrounds:
        if bg:
            if hasattr(bg, 'rgb'):
                r, g, b = bg.rgb
            else:
                r, g, b = bg
            avg = (r + g + b) / 3
            if avg < 80:
                theme = "dark"
            elif avg > 180:
                theme = "light"
            else:
                theme = "medium"
            break
    
    print(f"SLIDE {slide_num}: THEME={theme}")
    if backgrounds:
        print(f"Backgrounds: {backgrounds}")