"""
Convert presentation.html to PowerPoint.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Brand colours
EPI_BLUE = RGBColor(0x0B, 0x2B, 0x46)
EPI_BLUE_DEEP = RGBColor(0x06, 0x1A, 0x2E)
EPI_RED = RGBColor(0xD3, 0x20, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF3, 0xF4, 0xF6)
TEXT_DARK = RGBColor(0x1F, 0x29, 0x37)
GREY_500 = RGBColor(0x6B, 0x72, 0x80)
GREY_400 = RGBColor(0x9C, 0xA3, 0xAF)
GREY_300 = RGBColor(0xD1, 0xD5, 0xDB)
SLATE_300 = RGBColor(0xCB, 0xD5, 0xE1)
CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)

W = Inches(13.33)
H = Inches(7.5)
PAD_X = Inches(1.2)
PAD_Y = Inches(0.7)
INNER_W = W - 2 * PAD_X


prs = Presentation()
prs.slide_width = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(blank_layout)


def bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def txbox(slide, left, top, width, height, text, font_size=18, bold=False,
          color=TEXT_DARK, align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tf_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tf_box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf_box


def heading(slide, text, top, dark=True, size=36):
    color = WHITE if dark else EPI_BLUE
    box = slide.shapes.add_textbox(PAD_X, top, INNER_W, Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color
    return box


def kicker(slide, text, top, dark=True):
    color = GREY_400 if dark else GREY_500
    txbox(slide, PAD_X, top, INNER_W, Inches(0.35), text.upper(),
          font_size=10, bold=True, color=color)


def body(slide, text, left, top, width, height, dark=True, size=13):
    color = GREY_300 if dark else GREY_500
    txbox(slide, left, top, width, height, text, font_size=size, color=color)


def red_bar(slide, left, top, height=Inches(0.04)):
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, Inches(0.06), height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = EPI_RED
    bar.line.fill.background()


def card(slide, left, top, width, height, title, text,
         dark=False, emphasis=False):
    # background rect
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if emphasis:
        shape.fill.solid()
        shape.fill.fore_color.rgb = EPI_BLUE
    elif dark:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x0D, 0x22, 0x35)
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = EPI_RED
    shape.line.width = Pt(1.5)

    # red left accent bar
    red_bar(slide, left, top, height)

    title_color = WHITE if (dark or emphasis) else EPI_BLUE
    text_color = GREY_300 if (dark or emphasis) else GREY_500

    inner_l = left + Inches(0.18)
    inner_w = width - Inches(0.25)

    txbox(slide, inner_l, top + Inches(0.18), inner_w, Inches(0.35),
          title, font_size=13, bold=True, color=title_color)
    txbox(slide, inner_l, top + Inches(0.55), inner_w, height - Inches(0.65),
          text, font_size=11, color=text_color)


def metric_card(slide, left, top, width, height, metric, label, text, dark=True):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0D, 0x22, 0x35) if dark else CARD_BG
    shape.line.color.rgb = EPI_RED
    shape.line.width = Pt(1.5)
    red_bar(slide, left, top, height)

    inner_l = left + Inches(0.18)
    inner_w = width - Inches(0.25)

    txbox(slide, inner_l, top + Inches(0.12), inner_w, Inches(0.55),
          metric, font_size=28, bold=True, color=WHITE if dark else EPI_BLUE)
    txbox(slide, inner_l, top + Inches(0.68), inner_w, Inches(0.28),
          label.upper(), font_size=9, bold=True, color=GREY_400)
    txbox(slide, inner_l, top + Inches(1.0), inner_w, height - Inches(1.1),
          text, font_size=10, color=GREY_300 if dark else GREY_500)


def step_pipeline(slide, steps, top, dark=True):
    """Draw a horizontal pipeline of steps with arrows."""
    n = len(steps)
    step_w = Inches(1.5)
    arrow_w = Inches(0.4)
    total = n * step_w + (n - 1) * arrow_w
    left_start = (W - total) / 2

    for i, step in enumerate(steps):
        lx = left_start + i * (step_w + arrow_w)
        shape = slide.shapes.add_shape(1, lx, top, step_w, Inches(0.55))
        shape.fill.solid()
        shape.fill.fore_color.rgb = EPI_RED
        shape.line.fill.background()
        txbox(slide, lx, top + Inches(0.1), step_w, Inches(0.4),
              step, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < n - 1:
            ax = lx + step_w
            txbox(slide, ax, top + Inches(0.1), arrow_w, Inches(0.4),
                  "→", font_size=16, bold=True,
                  color=EPI_RED if not dark else SLATE_300,
                  align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────
s = add_slide()
bg(s, EPI_BLUE_DEEP)
heading(s, "Scaling AMS Through Intelligent Automation", Inches(1.8), dark=True, size=48)
txbox(s, PAD_X, Inches(3.4), INNER_W, Inches(0.9),
      "Strengthen delivery, improve margin, and position EPI-USE Africa as the premier AI-enabled SuccessFactors support partner in the region.",
      font_size=16, color=SLATE_300)
txbox(s, PAD_X, Inches(5.0), INNER_W, Inches(0.4),
      "Marco Smith  ·  EPI-USE Africa  ·  2026", font_size=12, color=GREY_400)

# ─────────────────────────────────────────────
# SLIDE 2 — AMS trust / unfair advantage
# ─────────────────────────────────────────────
s = add_slide()
bg(s, EPI_BLUE_DEEP)
heading(s,
        "AMS is where we earn our clients' trust.\nIn the AI era, that foundation becomes our unfair advantage.",
        Inches(1.5), dark=True, size=34)
txbox(s, PAD_X, Inches(4.2), INNER_W, Inches(1.8),
      "SAP is moving from assistant-led experiences to agentic, workflow-level execution — whether clients are ready or not.\n\n"
      "The opportunity is to lead that shift for our clients before the market forces it on them.",
      font_size=15, color=SLATE_300)

# ─────────────────────────────────────────────
# SLIDE 3 — About Me
# ─────────────────────────────────────────────
s = add_slide()
bg(s, WHITE)
kicker(s, "Introduction & About Me", PAD_Y, dark=False)
heading(s, "Why I am positioned to lead this now", PAD_Y + Inches(0.4), dark=False, size=30)

col_w = INNER_W / 2 - Inches(0.15)
card(s, PAD_X, Inches(2.0), col_w, Inches(3.8),
     "Current role",
     "Managing client relationships, escalations, contracts, and operational excellence across complex AMS environments.\n\n"
     "• Driving delivery quality in live, high-pressure settings\n"
     "• Bridging service, commercials, and client trust\n"
     "• Translating delivery pain points into product opportunities",
     dark=False)

card(s, PAD_X + col_w + Inches(0.3), Inches(2.0), col_w, Inches(3.8),
     "Key clients",
     "Trusted across demanding environments including:\n\n"
     "Allan Gray  ·  Discovery  ·  Capitec  ·  Thungela  ·  Mediclinic  ·  O&L  ·  Sanlam  ·  Vitality  ·  Illovo  ·  Shutterfly  ·  SARS  ·  UK Clients\n\n"
     "Notable: AGA  ·  UCT  ·  Ricoh  ·  Liberty  ·  MVAF  ·  SA Mint  ·  Engen",
     dark=False)

# ─────────────────────────────────────────────
# SLIDE 4 — Track Record
# ─────────────────────────────────────────────
s = add_slide()
bg(s, WHITE)
kicker(s, "Track Record & Impact", PAD_Y, dark=False)
heading(s, "Earned trust across complex, regulated, and demanding environments",
        PAD_Y + Inches(0.4), dark=False, size=28)

cw = INNER_W / 3 - Inches(0.15)
for i, (title, text) in enumerate([
    ("Operational excellence",
     "Consistently high SLA adherence, strong renewal momentum, and meaningful escalation reduction across managed clients."),
    ("Commercial impact",
     "Deep involvement in proposals, complex contract negotiations, and identifying upsell and white-space opportunities."),
    ("Leadership in action",
     "Already operating across delivery, commercials, relationships, and innovation at Associate Manager level."),
]):
    lx = PAD_X + i * (cw + Inches(0.2))
    card(s, lx, Inches(2.0), cw, Inches(1.7), title, text, dark=False)

for i, (title, text) in enumerate([
    ("Discovery", "Earned trust through consistent delivery in a demanding enterprise environment."),
    ("Capitec", "Demonstrated ability to manage delivery in a complex, regulated environment."),
    ("Allan Gray", "Proven track record of managing escalations, driving operational excellence, and supporting commercial growth."),
]):
    lx = PAD_X + i * (cw + Inches(0.2))
    card(s, lx, Inches(3.9), cw, Inches(1.5), title, text, dark=False)

# ─────────────────────────────────────────────
# SLIDE 5 — AI Landscape
# ─────────────────────────────────────────────
s = add_slide()
bg(s, EPI_BLUE_DEEP)
kicker(s, "The AI Landscape", PAD_Y, dark=True)
heading(s, "Understanding the yesterday, today, and tomorrow of AI in service delivery",
        PAD_Y + Inches(0.4), dark=True, size=28)

cw = INNER_W / 3 - Inches(0.15)
data = [
    ("Yesterday", "The chat era",
     "AI as a conversational assistant: useful for research, drafting, and troubleshooting, but still largely passive."),
    ("Today", "The agentic shift",
     "AI as an active participant: reading tickets, logs, and documentation to classify, diagnose, summarise, and recommend next actions."),
    ("Tomorrow", "Intelligent operations",
     "AI as an operator within guardrails: testing, validating, correcting, documenting, and escalating only the true edge cases."),
]
for i, (period, title, text) in enumerate(data):
    lx = PAD_X + i * (cw + Inches(0.2))
    top = Inches(2.2)
    h = Inches(3.6)
    card(s, lx, top, cw, h, f"{period.upper()} — {title}", text, dark=True)

txbox(s, PAD_X, Inches(6.2), INNER_W, Inches(0.5),
      "The shift is clear: AI is moving from  assistant → actor → operator",
      font_size=13, color=SLATE_300)

# ─────────────────────────────────────────────
# SLIDE 6 — SAP's 2026 AI Roadmap
# ─────────────────────────────────────────────
s = add_slide()
bg(s, EPI_BLUE_DEEP)
kicker(s, "SAP's 2026 AI Roadmap", PAD_Y, dark=True)
heading(s, "The ecosystem is already moving to a suite-wide agentic model",
        PAD_Y + Inches(0.4), dark=True, size=30)

txbox(s, PAD_X, Inches(1.85), INNER_W, Inches(0.5),
      "SAP's 1H 2026 SuccessFactors release expands Joule across recruiting, workforce administration, payroll, learning, performance, and talent development.",
      font_size=12, color=SLATE_300)

cw = INNER_W / 2 - Inches(0.15)
ch = Inches(1.8)
cards4 = [
    ("Joule's evolution",
     "The direction is no longer experimental. SAP is positioning agentic AI as a suite-level operating capability."),
    ("Scale of change",
     "As of Q1 2026, SAP reported more than 40 specialised agents and over 2,400 Joule Skills."),
    ("Specialised HR agents",
     "Career and Talent Development, HR Service, Payroll, and People Intelligence use cases are already being framed into real workflows."),
    ("What that means for us",
     "SAP is moving now. The only real question is whether our clients will be ready."),
]
for i, (title, text) in enumerate(cards4):
    row, col = divmod(i, 2)
    lx = PAD_X + col * (cw + Inches(0.3))
    top = Inches(2.5) + row * (ch + Inches(0.2))
    card(s, lx, top, cw, ch, title, text, dark=True)

# urgent box
ub = slide_shape = s.shapes.add_shape(1, PAD_X, Inches(6.3), INNER_W, Inches(0.5))
ub.fill.solid()
ub.fill.fore_color.rgb = EPI_RED
ub.line.fill.background()
txbox(s, PAD_X + Inches(0.1), Inches(6.32), INNER_W, Inches(0.45),
      "Readiness is becoming the differentiator.", font_size=13, bold=True, color=WHITE)

# ─────────────────────────────────────────────
# SLIDE 7 — AMS AI-Readiness
# ─────────────────────────────────────────────
s = add_slide()
bg(s, EPI_BLUE_DEEP)
kicker(s, "The AMS AI-Readiness Opportunity", PAD_Y, dark=True)
heading(s, "AI tools are only as good as the data, permissions, and configuration beneath them",
        PAD_Y + Inches(0.4), dark=True, size=28)

txbox(s, PAD_X, Inches(2.0), INNER_W, Inches(0.5),
      "In the AI era, AMS becomes the function that makes automation usable, trusted, and scalable.",
      font_size=14, bold=True, color=WHITE)

cw = INNER_W / 3 - Inches(0.15)
for i, (title, text) in enumerate([
    ("The bottleneck",
     "Poor data quality, weak job and position structures, outdated time configuration, and misaligned role-based permissions will undermine SAP's new AI capabilities."),
    ("The AMS mandate",
     "We must proactively prepare client environments now so they can switch Joule on with confidence tomorrow."),
    ("Strategic implication",
     "AI does not reduce the importance of AMS. It makes AMS more valuable because readiness becomes the new barrier to adoption."),
]):
    lx = PAD_X + i * (cw + Inches(0.2))
    card(s, lx, Inches(2.7), cw, Inches(3.5), title, text, dark=True)

# ─────────────────────────────────────────────
# SLIDE 8 — The Hard Metrics
# ─────────────────────────────────────────────
s = add_slide()
bg(s, EPI_BLUE_DEEP)
kicker(s, "The Hard Metrics", PAD_Y, dark=True)
heading(s, "AI is not just a technology story. It is an efficiency story.",
        PAD_Y + Inches(0.4), dark=True, size=32)

cw = INNER_W / 2 - Inches(0.15)
ch = Inches(2.4)
metrics = [
    ("4.87 hrs", "Saved per incident",
     "GenAI summarises tickets, searches knowledge bases, drafts responses, and pulls relevant logs — cutting 4.87 hours per incident."),
    ("30.5%", "Faster resolution",
     "AI agents auto-classify, prioritise, and route tickets — letting humans focus on diagnosis, not triage."),
    ("65.7%", "Ticket deflection",
     "Self-service AI answering password resets, policy questions, and simple requests — reducing human-agent demand by nearly two-thirds."),
    ("20%+", "Lower resolution cost",
     "Lower-cost AI agents handling L1-L2 issues, with humans escalating only for complex diagnosis and client interaction."),
]
for i, (metric, label, text) in enumerate(metrics):
    row, col = divmod(i, 2)
    lx = PAD_X + col * (cw + Inches(0.3))
    top = Inches(2.0) + row * (ch + Inches(0.15))
    metric_card(s, lx, top, cw, ch, metric, label, text, dark=True)

txbox(s, PAD_X, Inches(6.95), INNER_W, Inches(0.4),
      "Intelligent automation changes cost-to-serve, speed-to-resolution, and the competitiveness of the delivery model.",
      font_size=11, color=GREY_400)

# ─────────────────────────────────────────────
# SLIDE 9 — Margin & Value Shift
# ─────────────────────────────────────────────
s = add_slide()
bg(s, WHITE)
kicker(s, "The Margin & Value Shift", PAD_Y, dark=False)
heading(s, "AMS moves from reactive support to an intelligent service partner",
        PAD_Y + Inches(0.4), dark=False, size=30)

cw = INNER_W / 3 - Inches(0.15)
for i, (title, text) in enumerate([
    ("Lower cost to serve",
     "Automation removes thousands of consultant hours spent on repetitive triage and shortens the path to resolution."),
    ("Margin growth",
     "We can deliver the same or better outcomes with leaner, more focused effort and better operating leverage."),
    ("Reinvesting time",
     "Freed capacity moves into advisory, optimisation, architecture, and higher-value client work."),
]):
    lx = PAD_X + i * (cw + Inches(0.2))
    card(s, lx, Inches(2.2), cw, Inches(2.5), title, text, dark=False)

txbox(s, PAD_X, Inches(5.0), INNER_W, Inches(0.8),
      "This is the evolution: AMS becomes less break-fix, more proactive, more strategic, and more commercially valuable.",
      font_size=13, color=GREY_500)

# ─────────────────────────────────────────────
# SLIDE 10 — JobFit PoC
# ─────────────────────────────────────────────
s = add_slide()
bg(s, WHITE)
kicker(s, "JobFit: The Proof of Concept", PAD_Y, dark=False)
heading(s, "This strategy is not theoretical. I have already started building it.",
        PAD_Y + Inches(0.4), dark=False, size=30)

col_w = INNER_W / 2 - Inches(0.15)
card(s, PAD_X, Inches(2.0), col_w, Inches(2.1),
     "What JobFit is",
     "An AI-powered CV screening SaaS designed for unbiased, compliant, and rapid screening of applicants.\n\n"
     "• Built in record time using AI\n• Delivers value through AI at its core\n• Designed as recurring SaaS",
     emphasis=True)
card(s, PAD_X + col_w + Inches(0.3), Inches(2.0), col_w, Inches(2.1),
     "Why it matters",
     "JobFit proves a repeatable model: real client problem + rapid AI build + SaaS packaging = scalable revenue with minimal ongoing consultant involvement.",
     emphasis=True)

cw3 = INNER_W / 3 - Inches(0.15)
for i, (title, text) in enumerate([
    ("Prototype Timeline",
     "4 days from client ask to working demo — rapid iteration proves the model works."),
    ("GenAI Core",
     "Interprets CV data regardless of format — extracting structured fields from any document."),
    ("High Volume",
     "Processes hundreds of applications per role without proportionally increasing effort."),
]):
    lx = PAD_X + i * (cw3 + Inches(0.2))
    card(s, lx, Inches(4.3), cw3, Inches(1.6), title, text, dark=False)

txbox(s, PAD_X, Inches(6.1), INNER_W, Inches(0.5),
      "When we combine delivery insight with rapid prototyping, we can create funded, repeatable solutions.",
      font_size=12, color=GREY_500)

# ─────────────────────────────────────────────
# SLIDE 11 — Other Opportunities
# ─────────────────────────────────────────────
s = add_slide()
bg(s, WHITE)
kicker(s, "Other Opportunities", PAD_Y, dark=False)
heading(s, "Beyond JobFit: the AI-assisted opportunities I see",
        PAD_Y + Inches(0.4), dark=False, size=30)

cw = INNER_W / 3 - Inches(0.15)
opps = [
    ("System Health Check & Monitor",
     "Proactive review of admin alerts and check tool errors before they become incidents or escalations."),
    ("Data Integrity Engine",
     "Automated validation of org structure, position data, employee and employment information across the instance."),
    ("SuccessFactors Benefits Integration",
     "Medical aid and life insurance package integration with South African vendors directly in SuccessFactors."),
    ("Business Rule Review",
     "Review of errors, inconsistencies, and old config that accumulate over time in Workflow, RBPs, and Pay Structures."),
    ("Unit Testing & Regression Bot",
     "Browser-use agent for performing repeatable regression testing after releases and configuration changes."),
    ("Support Agent",
     "Lower tier support agent for resolving simple user issues or questions, working alongside Joule."),
]
for i, (title, text) in enumerate(opps):
    row, col = divmod(i, 3)
    lx = PAD_X + col * (cw + Inches(0.2))
    top = Inches(2.0) + row * Inches(2.0)
    card(s, lx, top, cw, Inches(1.8), title, text, dark=False)

txbox(s, PAD_X, Inches(6.2), INNER_W, Inches(0.4),
      "Each follows the Prototype Factory model: client problem → prototype → validate → fund → productise → scale.",
      font_size=12, color=GREY_500)

# ─────────────────────────────────────────────
# SLIDE 12 — Prototype Factory
# ─────────────────────────────────────────────
s = add_slide()
bg(s, EPI_BLUE_DEEP)
kicker(s, "The Prototype Factory", PAD_Y, dark=True)
heading(s, "A repeatable engine for funded innovation",
        PAD_Y + Inches(0.4), dark=True, size=34)

txbox(s, PAD_X, Inches(2.0), INNER_W, Inches(0.5),
      "Identify an AMS client problem, build a prototype in one week, validate early, secure buy-in, then productise and scale.",
      font_size=13, color=SLATE_300)

step_pipeline(s, ["Identify", "Prototype", "Validate", "Fund", "Productise", "Scale"],
              Inches(3.0), dark=True)

cw = INNER_W / 2 - Inches(0.15)
card(s, PAD_X, Inches(4.2), cw, Inches(1.5),
     "Why it works",
     "Client validation happens before major development investment.",
     dark=True)

# ─────────────────────────────────────────────
# SLIDE 13 — The Flywheel (Strategic Engine)
# ─────────────────────────────────────────────
s = add_slide()
bg(s, EPI_BLUE_DEEP)
kicker(s, "The Strategic Engine", PAD_Y, dark=True)
heading(s, "AMS is the flywheel that powers the whole model",
        PAD_Y + Inches(0.4), dark=True, size=34)

# Draw a simple text-based flywheel representation
# Centre circle
cx = W / 2
cy = Inches(4.0)
r = Inches(1.2)

# Three node circles at 120° intervals
import math
nodes = [
    ("Listen", "Client insight", -90),    # top
    ("Build", "Co-create", 30),           # bottom-right
    ("Scale", "Distribute", 150),         # bottom-left
]
for label, sub, deg in nodes:
    rad = math.radians(deg)
    nx = cx + r * math.cos(rad)
    ny = cy + r * math.sin(rad)
    node_r = Inches(0.55)
    shape = s.shapes.add_shape(9, nx - node_r, ny - node_r, node_r * 2, node_r * 2)  # 9 = oval
    shape.fill.solid()
    shape.fill.fore_color.rgb = EPI_BLUE
    shape.line.color.rgb = EPI_RED
    shape.line.width = Pt(2)
    txbox(s, nx - node_r, ny - node_r * 0.6, node_r * 2, Inches(0.35),
          label, font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, nx - node_r, ny + Inches(0.0), node_r * 2, Inches(0.25),
          sub, font_size=9, color=GREY_400, align=PP_ALIGN.CENTER)

# Flywheel label in centre
txbox(s, cx - Inches(0.8), cy - Inches(0.25), Inches(1.6), Inches(0.5),
      "AMS\nFLYWHEEL", font_size=10, color=GREY_400, align=PP_ALIGN.CENTER)

# Cards on the right side
card_x = cx + r + Inches(0.8)
card_w = W - card_x - PAD_X
ch = Inches(1.4)
for i, (title, text) in enumerate([
    ("Listen",
     "AMS relationships surface real client pain points first-hand — earlier and more honestly than any market research."),
    ("Build",
     "We co-create solutions with those clients using the Prototype Factory — validating real problems before any major investment."),
    ("Scale",
     "The AMS network becomes our lowest-friction distribution channel — existing trust lets us sell across the entire portfolio."),
]):
    card(s, card_x, Inches(2.0) + i * (ch + Inches(0.15)), card_w, ch, title, text, dark=True)

txbox(s, PAD_X, Inches(6.9), INNER_W, Inches(0.5),
      "Each cycle deepens client trust, sharpens insight, and widens the competitive moat. No competitor without an AMS network can replicate this.",
      font_size=11, color=SLATE_300)

# ─────────────────────────────────────────────
# SLIDE 14 — White Space Opportunity
# ─────────────────────────────────────────────
s = add_slide()
bg(s, EPI_BLUE_DEEP)
kicker(s, "The White Space Opportunity", PAD_Y, dark=True)
heading(s, "AMS is our lowest-friction distribution channel for innovation",
        PAD_Y + Inches(0.4), dark=True, size=30)

cw = INNER_W / 2 - Inches(0.15)
ch = Inches(2.0)
ws_cards = [
    ("Module gaps",
     "There is major white space across existing clients in Performance & Goals, Compensation, Learning, and Work Zone."),
    ("The captive market",
     "Every existing AMS client is already a warm lead for new AI-enabled solutions."),
    ("Our unfair advantage",
     "We already own the relationships, understand the pain points, and see the patterns earlier than anyone else."),
    ("Market positioning",
     "This is how we secure first-mover advantage as the premier AI-enabled SuccessFactors support partner in the region."),
]
for i, (title, text) in enumerate(ws_cards):
    row, col = divmod(i, 2)
    lx = PAD_X + col * (cw + Inches(0.3))
    top = Inches(2.2) + row * (ch + Inches(0.2))
    card(s, lx, top, cw, ch, title, text, dark=True)

# ─────────────────────────────────────────────
# SLIDE 15 — Risks & Mitigations
# ─────────────────────────────────────────────
s = add_slide()
bg(s, EPI_BLUE_DEEP)
kicker(s, "Risks & Mitigations", PAD_Y, dark=True)
heading(s, "The risks are real, but they are manageable by design",
        PAD_Y + Inches(0.4), dark=True, size=32)

cw = INNER_W / 2 - Inches(0.15)
ch = Inches(2.0)
risks = [
    ("Capacity constraints",
     "Cross-training and targeted hiring ensure AMS delivery quality never dips while we innovate."),
    ("Data privacy (POPIA)",
     "Responsible-by-design controls must be built into every prototype from day one, especially for high-compliance clients."),
    ("SaaS adoption",
     "Pilot programmes and internal client champions can create early traction before wider roll-out."),
    ("Consultant resistance",
     "Clear upskilling paths position AI as a career enabler and force multiplier, not a replacement."),
]
for i, (title, text) in enumerate(risks):
    row, col = divmod(i, 2)
    lx = PAD_X + col * (cw + Inches(0.3))
    top = Inches(2.2) + row * (ch + Inches(0.2))
    card(s, lx, top, cw, ch, title, text, dark=True)

# ─────────────────────────────────────────────
# SLIDE 16 — The Ask
# ─────────────────────────────────────────────
s = add_slide()
bg(s, WHITE)
kicker(s, "The Ask", PAD_Y, dark=False)
heading(s, "What I need to scale this properly",
        PAD_Y + Inches(0.4), dark=False, size=34)

asks = [
    ("1", "Guidance:",
     "A seat at the table in leadership and executive client discussions so I can hear challenges first-hand and identify prototyping opportunities."),
    ("2", "Assistance:",
     "Leadership support in positioning and selling these innovation solutions to existing client executives."),
    ("3", "Endorsement:",
     "Formal backing for the Prototype Factory as an official strategic initiative."),
    ("4", "Resourcing:",
     "Approval to cross-train, dedicate capacity to product development, and hire where needed."),
]

for i, (num, label, text) in enumerate(asks):
    top = Inches(2.1) + i * Inches(1.1)
    # number circle
    circ = s.shapes.add_shape(9, PAD_X, top + Inches(0.1), Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = EPI_RED
    circ.line.fill.background()
    txbox(s, PAD_X, top + Inches(0.1), Inches(0.5), Inches(0.5),
          num, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, PAD_X + Inches(0.65), top, INNER_W - Inches(0.65), Inches(0.35),
          label, font_size=13, bold=True, color=EPI_BLUE)
    txbox(s, PAD_X + Inches(0.65), top + Inches(0.35), INNER_W - Inches(0.65), Inches(0.65),
          text, font_size=12, color=GREY_500)

# ─────────────────────────────────────────────
# SLIDE 17 — Closing quote
# ─────────────────────────────────────────────
s = add_slide()
bg(s, EPI_BLUE_DEEP)

box = s.shapes.add_textbox(PAD_X, Inches(1.8), INNER_W, Inches(3.0))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "Back a model that is "
run.font.size = Pt(32)
run.font.bold = True
run.font.color.rgb = WHITE
run2 = p.add_run()
run2.text = "already running"
run2.font.size = Pt(32)
run2.font.bold = True
run2.font.color.rgb = EPI_RED

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.LEFT
r = p2.add_run()
r.text = ", scale it across the portfolio, and let me lead the team that delivers it."
r.font.size = Pt(32)
r.font.bold = True
r.font.color.rgb = WHITE

txbox(s, PAD_X, Inches(5.0), INNER_W, Inches(1.5),
      "AMS is the firm's closest, most consistent touchpoint with clients. That proximity is our best source of insight, "
      "our best route to adoption, and our best platform for AI-led growth.",
      font_size=15, color=SLATE_300)

# ─────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────
out_path = "presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}  ({len(prs.slides)} slides)")
